from __future__ import annotations

from collections import Counter
import os
import random
import re
from pathlib import Path

try:
    from pypdf import PdfReader
except Exception:  # pragma: no cover - optional dependency in some environments
    PdfReader = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional dependency in some environments
    Presentation = None


STOPWORDS = {
    "a",
    "about",
    "after",
    "all",
    "also",
    "an",
    "and",
    "are",
    "as",
    "at",
    "be",
    "because",
    "been",
    "before",
    "between",
    "but",
    "by",
    "can",
    "could",
    "did",
    "do",
    "does",
    "during",
    "each",
    "for",
    "from",
    "had",
    "has",
    "have",
    "he",
    "her",
    "his",
    "how",
    "i",
    "if",
    "in",
    "into",
    "is",
    "it",
    "its",
    "just",
    "may",
    "more",
    "most",
    "must",
    "not",
    "of",
    "on",
    "or",
    "our",
    "out",
    "over",
    "should",
    "so",
    "some",
    "such",
    "than",
    "that",
    "the",
    "their",
    "them",
    "then",
    "there",
    "these",
    "they",
    "this",
    "those",
    "through",
    "to",
    "under",
    "up",
    "use",
    "very",
    "was",
    "we",
    "were",
    "what",
    "when",
    "which",
    "who",
    "will",
    "with",
    "would",
    "you",
    "your",
}


def slugify(value: str) -> str:
    value = value.lower().strip()
    value = re.sub(r"[^a-z0-9]+", "-", value)
    return value.strip("-")


def split_sentences(text: str) -> list[str]:
    parts = re.split(r"(?<=[.!?])\s+", text.strip())
    return [part.strip() for part in parts if part.strip()]


def tokenize(text: str) -> list[str]:
    return [token for token in re.findall(r"[a-zA-Z0-9']+", text.lower()) if token not in STOPWORDS]


def top_terms(text: str, limit: int = 12) -> list[str]:
    terms = Counter(tokenize(text))
    return [term for term, _ in terms.most_common(limit)]


def extractive_summary(text: str, limit: int = 5) -> list[str]:
    sentences = split_sentences(text)
    if not sentences:
        return []

    terms = Counter(tokenize(text))

    scored = []
    for index, sentence in enumerate(sentences):
        score = sum(terms[token] for token in tokenize(sentence))
        scored.append((score, -index, sentence))

    scored.sort(reverse=True)
    summary = []
    for _, __, sentence in scored:
        if sentence not in summary:
            summary.append(sentence)
        if len(summary) >= limit:
            break
    return summary


def chunked(items: list[str], size: int) -> list[list[str]]:
    return [items[index : index + size] for index in range(0, len(items), size)]


def extract_text_from_upload(file_path: str) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix in {".txt", ".md", ".csv"}:
        return path.read_text(encoding="utf-8", errors="ignore")

    if suffix == ".pdf":
        if PdfReader is None:
            return ""
        reader = PdfReader(str(path))
        text = []
        for page in reader.pages:
            text.append(page.extract_text() or "")
        return "\n".join(text)

    if suffix == ".pptx":
        if Presentation is None:
            return ""
        presentation = Presentation(str(path))
        collected = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    collected.append(shape.text)
        return "\n".join(collected)

    if suffix in {".json", ".log"}:
        return path.read_text(encoding="utf-8", errors="ignore")

    return ""


def build_study_notes(text: str, subject_name: str | None = None) -> dict[str, object]:
    cleaned = re.sub(r"\s+", " ", text).strip()
    terms = top_terms(cleaned, limit=12)
    summary = extractive_summary(cleaned, limit=6)

    flashcards = []
    for term in terms[:6]:
        flashcards.append(
            {
                "front": f"What does {term} mean?",
                "back": f"Use the surrounding notes to define {term} clearly and add one exam-style example.",
            }
        )

    if subject_name and subject_name not in cleaned:
        summary.insert(0, f"Subject focus: {subject_name}")

    if not summary and cleaned:
        summary = [cleaned[:220] + ("..." if len(cleaned) > 220 else "")]

    if not terms and cleaned:
        terms = [part.strip(".,:;") for part in cleaned.split()[:8]]

    return {
        "summary": summary,
        "terms": terms,
        "flashcards": flashcards,
        "source_excerpt": cleaned[:900],
    }


def build_quiz_items(subject_name: str, topics: list[str], text: str = "") -> list[dict[str, object]]:
    base_terms = top_terms(text, limit=8)
    seeds = topics[:3] or base_terms[:3] or [subject_name.lower()]
    quiz_items = []

    for seed in seeds:
        distractors = [term for term in base_terms if term != seed][:2]
        if len(distractors) < 2:
            distractors.extend([f"{seed} example", f"{seed} rule"])
            distractors = distractors[:2]
        options = [seed, *distractors]
        random.shuffle(options)
        quiz_items.append(
            {
                "prompt": f"Which idea is most closely linked to {seed}?",
                "options": options,
                "answer_index": options.index(seed),
                "explanation": f"{seed} is a key study anchor for {subject_name}. Use it as a recall trigger.",
            }
        )
    return quiz_items


def build_revision_plan(subject_name: str, days: int = 7, weak_topics: list[str] | None = None) -> list[dict[str, str]]:
    weak_topics = weak_topics or []
    tasks = [
        "Quick recall sprint",
        "Notes review",
        "Topic drill",
        "Past-paper timed block",
        "Mark-scheme reflection",
        "Error log cleanup",
        "Mini mock",
    ]
    plan = []
    for day in range(1, days + 1):
        task = tasks[(day - 1) % len(tasks)]
        focus = weak_topics[(day - 1) % len(weak_topics)] if weak_topics else subject_name
        plan.append(
            {
                "day": f"Day {day}",
                "task": task,
                "focus": focus,
                "goal": "30-45 minutes",
            }
        )
    return plan


def filter_results(query: str, haystack: list[dict[str, str]]) -> list[dict[str, str]]:
    normalized = query.lower().strip()
    if not normalized:
        return haystack
    results = []
    for item in haystack:
        blob = " ".join(item.values()).lower()
        if normalized in blob:
            results.append(item)
    return results
